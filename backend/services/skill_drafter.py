"""
Document-to-Skill Drafter Service
Extracts text from PDF/PPTX/DOCX, detects topic, generates a draft skill via LLM.
"""

import io
import re
import json
import logging
from typing import Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from PDF bytes using pypdf (installed) or pdfminer fallback."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except ImportError:
        pass

    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        return pdfminer_extract(io.BytesIO(data))
    except ImportError:
        pass

    raise RuntimeError(
        "PDF extraction requires pypdf or pdfminer.six. "
        "Install via: pip install pypdf"
    )


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from DOCX bytes using python-docx."""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # also grab table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n\n".join(paragraphs)
    except ImportError:
        raise RuntimeError(
            "DOCX extraction requires python-docx. "
            "Install via: pip install python-docx"
        )


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from PPTX bytes using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append("\n".join(slide_texts))
        return "\n\n--- SLIDE BREAK ---\n\n".join(parts)
    except ImportError:
        raise RuntimeError(
            "PPTX extraction requires python-pptx. "
            "Install via: pip install python-pptx"
        )


def extract_text(filename: str, data: bytes) -> str:
    """Dispatch to the right extractor based on file extension."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(data)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(data)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(data)
    elif ext in (".txt", ".md"):
        return data.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: PDF, DOCX, PPTX, TXT, MD")


# ---------------------------------------------------------------------------
# Topic detection (heuristic, fast — no LLM)
# ---------------------------------------------------------------------------

TOPIC_PATTERNS = [
    # (regex pattern, topic_slug, topic_label)
    (r"thu[ế\s]+TNCN|thu nhập cá nhân|PIT|personal income tax", "pit", "Thuế TNCN (PIT)"),
    (r"thu[ế\s]+GTGT|thuế giá trị gia tăng|VAT|GTGT", "vat", "Thuế GTGT (VAT)"),
    (r"thu[ế\s]+TNDN|thu nhập doanh nghiệp|CIT|corporate income", "cit", "Thuế TNDN (CIT)"),
    (r"hải quan|customs|xuất nhập khẩu|import.*export|HS code|HS\s*\d{4}", "customs", "Hải quan (Customs)"),
    (r"chuyển giá|transfer pricing|giá liên kết|arm.s length", "transfer-pricing", "Chuyển giá (Transfer Pricing)"),
    (r"nhà thầu nước ngoài|FCT|foreign contractor|withholding", "fct", "Thuế nhà thầu (FCT)"),
    (r"tiêu thụ đặc biệt|SCT|special consumption|excise", "sct", "Thuế TTĐB (SCT)"),
    (r"quản lý thuế|tax administration|thuế điện tử|eTax|kê khai", "tax-admin", "Quản lý thuế"),
    (r"BHXH|bảo hiểm xã hội|social insurance|BHYT|BHTN", "social-insurance", "Bảo hiểm xã hội"),
    (r"hiệp định thuế|DTA|double tax|tax treaty", "dta", "Hiệp định thuế (DTA)"),
]


def detect_topic(text: str) -> tuple[str, str]:
    """Returns (topic_slug, topic_label) based on keyword frequency."""
    text_lower = text.lower()
    scores: dict[str, int] = {}
    labels: dict[str, str] = {}

    for pattern, slug, label in TOPIC_PATTERNS:
        matches = len(re.findall(pattern, text, re.IGNORECASE))
        if matches > 0:
            scores[slug] = matches
            labels[slug] = label

    if not scores:
        return "general", "Tư vấn thuế chung"

    best_slug = max(scores, key=lambda k: scores[k])
    return best_slug, labels[best_slug]


# ---------------------------------------------------------------------------
# LLM-powered draft skill generator
# ---------------------------------------------------------------------------

DRAFT_SYSTEM_PROMPT = """You are a Vietnamese tax expert AI assistant. 
Your task is to convert a provided document (legal text, tax guidance, regulation, etc.) 
into a structured Skill file in Markdown format following the TaxLegal AI skill format.

The output must be a valid skill file with:
1. YAML frontmatter (between --- delimiters) with: name, version, description, category, tags, applicable_bots, editable
2. A well-structured Markdown body covering:
   - Key rules and regulations extracted from the document
   - Important numbers, rates, thresholds, deadlines
   - Practical guidance and examples where applicable
   - Anti-hallucination rules or caveats

IMPORTANT:
- Write primarily in Vietnamese (mixed with English technical terms as needed)
- Keep the content accurate to the source document — do not hallucinate
- Format tables, code blocks, and headers properly
- Set editable: true in frontmatter
- Set applicable_bots: [partner, ja, sa]
- Suggest appropriate category: tax | customs | compliance | advisory
- Keep the skill focused and concise (aim for 300-600 lines max)
"""


async def generate_draft_skill(
    filename: str,
    extracted_text: str,
    topic_slug: str,
    topic_label: str,
    ai_provider=None,  # kept for backward compat but ignored
    max_tokens: int = 4096,
) -> str:
    """
    Call the LLM to generate a draft skill from extracted document text.
    Returns raw markdown string (frontmatter + body).
    Uses call_ai() from backend.ai_provider directly.
    """
    from backend.ai_provider import call_ai

    # Truncate text if too long (keep first 8000 chars to fit context)
    if len(extracted_text) > 8000:
        truncated = extracted_text[:8000]
        truncation_notice = (
            f"\n\n[NOTE: Document truncated to 8000 chars. Original length: {len(extracted_text)} chars.]"
        )
    else:
        truncated = extracted_text
        truncation_notice = ""

    user_prompt = f"""Convert the following document into a TaxLegal AI skill file.

Source document: {filename}
Detected topic: {topic_label} ({topic_slug})

Document content:
---
{truncated}{truncation_notice}
---

Generate the skill file now. Start with the YAML frontmatter (---) immediately.
Suggested skill name: {topic_slug}-from-document
"""

    try:
        response = await call_ai(
            model_id="deepseek-chat",
            messages=[{"role": "user", "content": user_prompt}],
            system_prompt=DRAFT_SYSTEM_PROMPT,
            max_tokens=max_tokens,
            temperature=0.3,
        )
        return response.get("content", "")
    except Exception as e:
        logger.error(f"LLM skill draft failed: {e}")
        # Return a minimal template with extracted text as fallback
        return _fallback_draft(filename, extracted_text, topic_slug, topic_label)


def _fallback_draft(
    filename: str, text: str, topic_slug: str, topic_label: str
) -> str:
    """Generate a minimal draft without LLM when LLM fails."""
    snippet = text[:2000].replace("---", "- - -")
    return f"""---
name: {topic_slug}-draft
version: 0.1.0
description: >
  Draft skill tự động tạo từ tài liệu: {filename}
  Chủ đề phát hiện: {topic_label}
  Cần review và chỉnh sửa trước khi sử dụng.
category: tax
tags: [{topic_slug}, draft, vietnam]
applicable_bots: [partner, ja, sa]
editable: true
---

# Skill Draft: {topic_label}

> **Nguồn tài liệu:** {filename}
> **Trạng thái:** DRAFT — Cần review và bổ sung trước khi dùng

---

## Nội Dung Tài Liệu Gốc (Trích dẫn)

```
{snippet}
```

---

## Ghi Chú

- Skill này được tạo tự động từ tài liệu nguồn
- Cần chuyên gia thuế review và hoàn chỉnh nội dung
- Các số liệu, mức thuế suất cần được xác nhận với văn bản pháp lý hiện hành
"""


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

async def create_draft_from_document(
    filename: str,
    data: bytes,
    ai_provider,
    use_llm: bool = True,
) -> dict:
    """
    Main entry point: extract text → detect topic → generate draft skill.
    
    Returns:
        {
            "filename": str,
            "extracted_text": str,
            "text_length": int,
            "topic_slug": str,
            "topic_label": str,
            "draft_content": str,  # full markdown skill file
            "suggested_slug": str,
            "error": str | None,
        }
    """
    result: dict = {
        "filename": filename,
        "extracted_text": "",
        "text_length": 0,
        "topic_slug": "general",
        "topic_label": "Tư vấn thuế chung",
        "draft_content": "",
        "suggested_slug": "",
        "error": None,
    }

    # Step 1: Extract text
    try:
        text = extract_text(filename, data)
        result["extracted_text"] = text
        result["text_length"] = len(text)
    except Exception as e:
        result["error"] = f"Text extraction failed: {e}"
        return result

    # Step 2: Detect topic
    topic_slug, topic_label = detect_topic(text)
    result["topic_slug"] = topic_slug
    result["topic_label"] = topic_label

    # Suggest a slug based on filename + topic
    safe_name = re.sub(r"[^a-z0-9]+", "-", Path(filename).stem.lower()).strip("-")
    result["suggested_slug"] = f"{topic_slug}-{safe_name}"[:60]

    # Step 3: Generate draft
    if use_llm:
        draft = await generate_draft_skill(
            filename=filename,
            extracted_text=text,
            topic_slug=topic_slug,
            topic_label=topic_label,
            ai_provider=ai_provider,
        )
    else:
        draft = _fallback_draft(filename, text, topic_slug, topic_label)

    result["draft_content"] = draft
    return result
